import streamlit as st
import os
import io
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from query_engine import get_query_engine, format_sources

# --- Chargement des variables d'environnement ---
load_dotenv()

if hasattr(st, "secrets"):
    for key, value in st.secrets.items():
        os.environ[key] = value

# --- Configuration de la page ---
st.set_page_config(
    page_title="ACENOS Knowledge Agent",
    page_icon="🏦",
    layout="wide"
)

# --- Header ---
st.title("🏦 ACENOS Knowledge Agent")
st.caption("Interroge la base de connaissances ACENOS · Sources citées automatiquement")

# --- Chargement de l'agent (mis en cache) ---
@st.cache_resource
def load_engine():
    return get_query_engine()

engine = load_engine()

# --- Sidebar ---
with st.sidebar:
    st.header("⚙️ Paramètres")

    # Format de sortie
    format_sortie = st.selectbox(
        "Format de réponse",
        ["Texte structuré", "Présentation PPT", "Mail professionnel"],
        help="Choisir le format dans lequel l'agent doit restituer sa réponse."
    )

    st.divider()

    # Prompt éditable
    with st.expander("✏️ Modifier le prompt système", expanded=False):
        prompt_path = os.path.join(os.path.dirname(__file__), "prompt.txt")
        try:
            with open(prompt_path, "r", encoding="utf-8") as f:
                prompt_actuel = f.read()
        except FileNotFoundError:
            prompt_actuel = "Fichier prompt.txt introuvable."

        nouveau_prompt = st.text_area(
            "Prompt système",
            value=prompt_actuel,
            height=350,
            help="Modifie les instructions données à l'agent. Clique sur Appliquer pour recharger."
        )
        if st.button("✅ Appliquer le prompt"):
            with open(prompt_path, "w", encoding="utf-8") as f:
                f.write(nouveau_prompt)
            st.cache_resource.clear()
            st.success("Prompt mis à jour — agent rechargé.")
            st.rerun()

    st.divider()

    # Bouton reset conversation
    if st.button("🗑️ Effacer la conversation"):
        st.session_state.messages = []
        st.rerun()

    st.caption("ACENOS Consulting · Knowledge Agent v1.0")

# --- Historique de conversation ---
if "messages" not in st.session_state:
    st.session_state.messages = []

# Message d'accueil si conversation vide
if not st.session_state.messages:
    with st.chat_message("assistant"):
        st.markdown(
            "Bonjour 👋 Je suis l'Agent de Connaissances ACENOS.\n\n"
            "Pose-moi une question sur nos missions, formations, retours d'expérience "
            "ou notre veille réglementaire. Je cite toujours mes sources.\n\n"
            "**Exemples :**\n"
            "- *Quelles sont nos meilleures pratiques pour une due diligence ESG en banque ?*\n"
            "- *Génère-moi les slides d'un pitch pour une mission de transformation comptable*\n"
            "- *Rédige un mail de présentation de nos offres DORA à un DAF bancaire*"
        )

# Affichage de l'historique
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# --- Fonction export PPT ---
def generer_ppt(question: str, contenu: str) -> bytes:
    """Génère un fichier PPTX à partir de la réponse de l'agent."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide de titre
    slide_titre = prs.slides.add_slide(prs.slide_layouts[0])
    slide_titre.shapes.title.text = question[:80]
    slide_titre.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    slide_titre.placeholders[1].text = "Base de connaissances ACENOS · Knowledge Agent"

    # Découpe le contenu en slides (1 slide par bloc séparé par double saut de ligne)
    blocs = [b.strip() for b in contenu.split("\n\n") if b.strip()]

    for bloc in blocs[:8]:  # max 8 slides de contenu
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lignes = bloc.split("\n")

        # Titre du slide = première ligne du bloc
        titre = lignes[0].lstrip("#").strip()[:80]
        slide.shapes.title.text = titre
        slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(24)

        # Corps = lignes suivantes
        corps = "\n".join(lignes[1:]).strip()
        if corps and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.text = corps[:800]
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(16)

    # Slide de conclusion
    slide_fin = prs.slides.add_slide(prs.slide_layouts[0])
    slide_fin.shapes.title.text = "Sources & Contacts"
    slide_fin.placeholders[1].text = "Document généré par ACENOS Knowledge Agent\ncontact@acenos.fr · www.acenos.fr"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- Zone de saisie principale ---
if question := st.chat_input("Pose ta question à la base de connaissances ACENOS..."):

    # Ajout de la question à l'historique
    st.session_state.messages.append({"role": "user", "content": question})
    with st.chat_message("user"):
        st.markdown(question)

    # Adaptation du prompt selon le format choisi
    if format_sortie == "Présentation PPT":
        prompt = (
            f"{question}\n\n"
            "Structure ta réponse comme un plan de présentation PowerPoint :\n"
            "- Un titre de slide par grande idée (précédé de ##)\n"
            "- 3 à 5 bullet points concis par slide\n"
            "- Maximum 8 slides\n"
            "- Commence par un slide 'Contexte & Enjeux' et termine par 'Recommandations'"
        )
    elif format_sortie == "Mail professionnel":
        prompt = (
            f"{question}\n\n"
            "Rédige ta réponse sous forme d'un mail professionnel prêt à envoyer :\n"
            "- Objet : [objet du mail]\n"
            "- Corps structuré avec introduction, développement et conclusion\n"
            "- Ton professionnel de consultant senior\n"
            "- Signature : [Prénom Nom] | ACENOS Consulting | contact@acenos.fr"
        )
    else:
        prompt = question

    # Appel à l'agent
    with st.chat_message("assistant"):
        with st.spinner("Recherche dans la base de connaissances ACENOS..."):
            try:
                response = engine.query(prompt)
                contenu_reponse = str(response)

                # Affichage de la réponse
                st.markdown(contenu_reponse)

                # Sources détaillées
                with st.expander("📎 Sources utilisées", expanded=True):
                    st.markdown(format_sources(response))

                # Bouton export PPT
                if format_sortie == "Présentation PPT":
                    ppt_bytes = generer_ppt(question, contenu_reponse)
                    st.download_button(
                        label="⬇️ Télécharger la présentation PowerPoint",
                        data=ppt_bytes,
                        file_name="acenos_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )

                # Ajout à l'historique
                st.session_state.messages.append({
                    "role": "assistant",
                    "content": contenu_reponse
                })

            except Exception as e:
                st.error(f"Erreur lors de la requête : {str(e)}")
                st.info("Vérifie que tes clés API sont bien configurées et que la base de connaissances est indexée.")